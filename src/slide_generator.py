"""
PowerPoint Slide Generator
Creates executive-ready one-page leadership slides from social media performance data.
"""

import io
from datetime import datetime
from typing import Dict, List, Optional, Tuple
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from .metrics_engine import PerformanceSnapshot, MetricResult
from .insights_engine import InsightsEngine


class SlideColors:
    """Color scheme for the leadership slide."""
    # VA Brand Colors (approximated)
    PRIMARY_BLUE = RGBColor(0, 51, 102)  # Dark blue
    ACCENT_BLUE = RGBColor(0, 94, 162)  # Medium blue
    LIGHT_BLUE = RGBColor(225, 243, 255)  # Light blue background
    
    # Status colors
    POSITIVE_GREEN = RGBColor(34, 139, 34)  # Forest green
    NEGATIVE_RED = RGBColor(178, 34, 34)  # Firebrick red
    NEUTRAL_GRAY = RGBColor(128, 128, 128)  # Gray
    
    # Text colors
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)
    DARK_GRAY = RGBColor(64, 64, 64)
    
    # Background
    CARD_BG = RGBColor(248, 249, 250)
    HEADER_BG = RGBColor(0, 51, 102)


class SlideGenerator:
    """
    Generates a one-page leadership slide from social media performance data.
    """
    
    def __init__(self, snapshot: PerformanceSnapshot, insights: Dict):
        self.snapshot = snapshot
        self.insights = insights
        self.prs = Presentation()
        
        # Set slide dimensions (16:9 widescreen)
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
        self.slide = None
    
    def generate_slide(self) -> Presentation:
        """
        Generate the complete leadership slide.
        
        Returns:
            PowerPoint Presentation object
        """
        # Add blank slide
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        self.slide = self.prs.slides.add_slide(slide_layout)
        
        # Build slide sections
        self._add_header()
        self._add_performance_snapshot()
        self._add_trend_highlights()
        self._add_top_content()
        self._add_insights_section()
        self._add_footer()
        
        return self.prs
    
    def _add_header(self):
        """Add header section with brand name, date range, and platforms."""
        # Header background
        header = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(0.9)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = SlideColors.HEADER_BG
        header.line.fill.background()
        
        # Title
        title_box = self.slide.shapes.add_textbox(
            Inches(0.3), Inches(0.15),
            Inches(6), Inches(0.4)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Social Media Performance Report"
        title_para.font.size = Pt(24)
        title_para.font.bold = True
        title_para.font.color.rgb = SlideColors.WHITE
        
        # Date range
        date_range = f"{self.snapshot.period_start.strftime('%b %d')} - {self.snapshot.period_end.strftime('%b %d, %Y')}"
        date_box = self.slide.shapes.add_textbox(
            Inches(0.3), Inches(0.55),
            Inches(4), Inches(0.3)
        )
        date_frame = date_box.text_frame
        date_para = date_frame.paragraphs[0]
        date_para.text = date_range
        date_para.font.size = Pt(14)
        date_para.font.color.rgb = SlideColors.WHITE
        
        # Platforms included
        platforms = list(self.snapshot.platform_breakdown.keys())
        if platforms:
            platform_text = " | ".join([p.title() for p in platforms])
            platform_box = self.slide.shapes.add_textbox(
                Inches(9), Inches(0.3),
                Inches(4), Inches(0.3)
            )
            platform_frame = platform_box.text_frame
            platform_para = platform_frame.paragraphs[0]
            platform_para.text = f"Platforms: {platform_text}"
            platform_para.font.size = Pt(12)
            platform_para.font.color.rgb = SlideColors.WHITE
            platform_para.alignment = PP_ALIGN.RIGHT
    
    def _add_performance_snapshot(self):
        """Add the key performance metrics section."""
        # Section title
        section_title = self.slide.shapes.add_textbox(
            Inches(0.3), Inches(1.0),
            Inches(4), Inches(0.35)
        )
        title_frame = section_title.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Performance Snapshot"
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = SlideColors.PRIMARY_BLUE
        
        # Get headline metrics
        headline_metrics = self.snapshot.get_headline_metrics()
        
        # Calculate card positions - 6 cards for more metrics
        start_x = 0.3
        card_width = 2.0
        card_height = 1.2
        card_spacing = 0.1
        start_y = 1.4
        
        for i, metric in enumerate(headline_metrics[:6]):
            x_pos = start_x + (i * (card_width + card_spacing))
            self._add_metric_card(x_pos, start_y, card_width, card_height, metric)
    
    def _add_metric_card(self, x: float, y: float, width: float, height: float, 
                         metric: MetricResult):
        """Add a single metric card."""
        # Card background
        card = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = SlideColors.CARD_BG
        card.line.color.rgb = RGBColor(220, 220, 220)
        card.line.width = Pt(1)
        
        # Metric value
        value_box = self.slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 0.15),
            Inches(width - 0.2), Inches(0.5)
        )
        value_frame = value_box.text_frame
        value_para = value_frame.paragraphs[0]
        value_para.text = metric.formatted_current
        value_para.font.size = Pt(28)
        value_para.font.bold = True
        value_para.font.color.rgb = SlideColors.PRIMARY_BLUE
        value_para.alignment = PP_ALIGN.CENTER
        
        # Change indicator
        change_box = self.slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 0.6),
            Inches(width - 0.2), Inches(0.25)
        )
        change_frame = change_box.text_frame
        change_para = change_frame.paragraphs[0]
        
        # Add trend arrow and percentage
        arrow = "▲" if metric.trend == 'up' else "▼" if metric.trend == 'down' else "―"
        change_para.text = f"{arrow} {metric.formatted_change}"
        change_para.font.size = Pt(14)
        change_para.font.bold = True
        
        if metric.is_positive_trend:
            change_para.font.color.rgb = SlideColors.POSITIVE_GREEN
        elif metric.trend == 'stable':
            change_para.font.color.rgb = SlideColors.NEUTRAL_GRAY
        else:
            change_para.font.color.rgb = SlideColors.NEGATIVE_RED
        
        change_para.alignment = PP_ALIGN.CENTER
        
        # Metric label
        label_box = self.slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 0.9),
            Inches(width - 0.2), Inches(0.25)
        )
        label_frame = label_box.text_frame
        label_para = label_frame.paragraphs[0]
        label_para.text = metric.display_name
        label_para.font.size = Pt(10)
        label_para.font.color.rgb = SlideColors.DARK_GRAY
        label_para.alignment = PP_ALIGN.CENTER
    
    def _add_trend_highlights(self):
        """Add trend highlights section."""
        # Section title
        section_title = self.slide.shapes.add_textbox(
            Inches(0.3), Inches(2.75),
            Inches(6), Inches(0.35)
        )
        title_frame = section_title.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "📈 Trend Highlights"
        title_para.font.size = Pt(16)
        title_para.font.bold = True
        title_para.font.color.rgb = SlideColors.PRIMARY_BLUE
        
        # Get trend highlights from insights engine
        engine = InsightsEngine(self.snapshot)
        highlights = engine.get_trend_highlights()
        
        # Add highlight items
        y_pos = 3.15
        for i, highlight in enumerate(highlights[:4]):
            self._add_trend_item(0.3, y_pos + (i * 0.4), highlight)
    
    def _add_trend_item(self, x: float, y: float, highlight: Dict):
        """Add a single trend highlight item."""
        # Trend indicator
        indicator = "▲" if highlight['trend'] == 'up' else "▼" if highlight['trend'] == 'down' else "―"
        color = SlideColors.POSITIVE_GREEN if highlight['is_positive'] else SlideColors.NEGATIVE_RED
        
        item_box = self.slide.shapes.add_textbox(
            Inches(x), Inches(y),
            Inches(6), Inches(0.35)
        )
        item_frame = item_box.text_frame
        item_para = item_frame.paragraphs[0]
        
        # Build text with formatting
        item_para.text = f"{indicator} {highlight['metric']}: {highlight['value']} ({highlight['change']})"
        item_para.font.size = Pt(12)
        item_para.font.color.rgb = color
    
    def _add_top_content(self):
        """Add top performing content section."""
        # Section background
        section_bg = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.6), Inches(2.75),
            Inches(6.4), Inches(2.0)
        )
        section_bg.fill.solid()
        section_bg.fill.fore_color.rgb = SlideColors.CARD_BG
        section_bg.line.color.rgb = RGBColor(220, 220, 220)
        
        # Section title
        section_title = self.slide.shapes.add_textbox(
            Inches(6.8), Inches(2.85),
            Inches(6), Inches(0.35)
        )
        title_frame = section_title.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "🏆 Top Performing Content"
        title_para.font.size = Pt(14)
        title_para.font.bold = True
        title_para.font.color.rgb = SlideColors.PRIMARY_BLUE
        
        # Top post details
        top_post = self.snapshot.top_post
        if top_post:
            # Post type badge
            type_box = self.slide.shapes.add_textbox(
                Inches(6.8), Inches(3.25),
                Inches(1.5), Inches(0.25)
            )
            type_frame = type_box.text_frame
            type_para = type_frame.paragraphs[0]
            type_para.text = f"[{top_post.get('post_type', 'Post')}]"
            type_para.font.size = Pt(10)
            type_para.font.bold = True
            type_para.font.color.rgb = SlideColors.ACCENT_BLUE
            
            # Post message preview
            message = top_post.get('message', 'No message available')
            if len(message) > 120:
                message = message[:120] + "..."
            
            msg_box = self.slide.shapes.add_textbox(
                Inches(6.8), Inches(3.5),
                Inches(6), Inches(0.7)
            )
            msg_frame = msg_box.text_frame
            msg_frame.word_wrap = True
            msg_para = msg_frame.paragraphs[0]
            msg_para.text = f'"{message}"'
            msg_para.font.size = Pt(10)
            msg_para.font.italic = True
            msg_para.font.color.rgb = SlideColors.DARK_GRAY
            
            # Key metrics
            engagement = top_post.get('engagement', 0)
            reach = top_post.get('reach', 0)
            
            metrics_box = self.slide.shapes.add_textbox(
                Inches(6.8), Inches(4.25),
                Inches(6), Inches(0.3)
            )
            metrics_frame = metrics_box.text_frame
            metrics_para = metrics_frame.paragraphs[0]
            metrics_para.text = f"Engagement: {engagement:,}  |  Reach: {reach:,}"
            metrics_para.font.size = Pt(11)
            metrics_para.font.bold = True
            metrics_para.font.color.rgb = SlideColors.PRIMARY_BLUE
    
    def _add_insights_section(self):
        """Add the insights and recommendations section."""
        # Section background
        section_bg = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.3), Inches(4.9),
            Inches(12.7), Inches(2.2)
        )
        section_bg.fill.solid()
        section_bg.fill.fore_color.rgb = RGBColor(245, 247, 250)
        section_bg.line.color.rgb = RGBColor(200, 200, 200)
        
        # Three columns: What Happened | Why It Matters | What To Do
        col_width = 4.0
        col_spacing = 0.15
        start_x = 0.5
        
        columns = [
            ("📋 What Happened", self.insights.get('what_happened', []), 0),
            ("💡 Why It Matters", self.insights.get('why_it_matters', []), 1),
            ("🎯 What To Do Next", self.insights.get('what_to_do', []), 2),
        ]
        
        for title, items, idx in columns:
            x_pos = start_x + (idx * (col_width + col_spacing))
            self._add_insight_column(x_pos, 5.0, col_width, title, items)
    
    def _add_insight_column(self, x: float, y: float, width: float, 
                           title: str, items: List[str]):
        """Add an insight column."""
        # Column title
        title_box = self.slide.shapes.add_textbox(
            Inches(x), Inches(y),
            Inches(width), Inches(0.35)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(12)
        title_para.font.bold = True
        title_para.font.color.rgb = SlideColors.PRIMARY_BLUE
        
        # Items
        items_box = self.slide.shapes.add_textbox(
            Inches(x), Inches(y + 0.4),
            Inches(width), Inches(1.6)
        )
        items_frame = items_box.text_frame
        items_frame.word_wrap = True
        
        for i, item in enumerate(items[:3]):
            if i == 0:
                para = items_frame.paragraphs[0]
            else:
                para = items_frame.add_paragraph()
            
            # Clean up emoji/bullet if already present
            clean_item = item.lstrip('✅⚠️📈📉🏆💡 ')
            para.text = f"• {clean_item}"
            para.font.size = Pt(10)
            para.font.color.rgb = SlideColors.DARK_GRAY
            para.space_after = Pt(6)
    
    def _add_footer(self):
        """Add footer with generation timestamp."""
        footer_box = self.slide.shapes.add_textbox(
            Inches(0.3), Inches(7.15),
            Inches(12.7), Inches(0.25)
        )
        footer_frame = footer_box.text_frame
        footer_para = footer_frame.paragraphs[0]
        footer_para.text = f"Generated on {datetime.now().strftime('%B %d, %Y at %I:%M %p')} | Campaign Analyzer"
        footer_para.font.size = Pt(8)
        footer_para.font.color.rgb = SlideColors.NEUTRAL_GRAY
        footer_para.alignment = PP_ALIGN.CENTER
    
    def save(self, filepath: str):
        """Save the presentation to a file."""
        self.prs.save(filepath)
    
    def to_bytes(self) -> bytes:
        """Return the presentation as bytes (for download)."""
        buffer = io.BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()


class PDFExporter:
    """
    Exports slides to PDF format.
    Note: Requires additional library or conversion service.
    """
    
    @staticmethod
    def export_to_pdf(pptx_bytes: bytes, output_path: str) -> bool:
        """
        Export PowerPoint to PDF.
        
        Note: This is a placeholder. Full PDF conversion typically requires
        either LibreOffice, unoconv, or a cloud service.
        """
        # For MVP, we'll save PPTX and note that PDF conversion
        # requires additional setup
        print("Note: PDF export requires LibreOffice or similar converter.")
        print("Saving as PPTX instead.")
        return False


def generate_leadership_slide(snapshot: PerformanceSnapshot, insights: Dict) -> SlideGenerator:
    """
    Convenience function to generate a leadership slide.
    
    Args:
        snapshot: Performance snapshot from MetricsCalculator
        insights: Insights dictionary from InsightsEngine
        
    Returns:
        SlideGenerator instance with generated slide
    """
    generator = SlideGenerator(snapshot, insights)
    generator.generate_slide()
    return generator
